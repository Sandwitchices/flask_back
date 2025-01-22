import os
import openai
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from werkzeug.utils import secure_filename
from pptx import Presentation
from docx import Document
import fitz  # PyMuPDF

app = Flask(__name__)

# Enable CORS for all domains
CORS(app)

# Configure upload folder
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['DOCX_FOLDER'] = 'docs'

# Ensure the upload and docs folders exist
if not os.path.exists(app.config['UPLOAD_FOLDER']):
    os.makedirs(app.config['UPLOAD_FOLDER'])
if not os.path.exists(app.config['DOCX_FOLDER']):
    os.makedirs(app.config['DOCX_FOLDER'])

# Set your OpenAI API key
openai.api_key = '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
'

def generate_summary(text):
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": f"Simplify and summarize the following content for students to review:\n\n{text}"}
        ],
        max_tokens=1024,
        n=1,
        stop=None,
        temperature=0.7,
    )
    summary = response['choices'][0]['message']['content'].strip()
    return summary

@app.route('/')
def home():
    return "Welcome to the PPT and PDF Parser API!"

@app.route('/parse', methods=['POST'])
def parse_file():
    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400

    filename = secure_filename(file.filename)
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(file_path)

    try:
        doc_content = ""
        if file.filename.endswith('.pptx'):
            # Parse the PowerPoint file
            presentation = Presentation(file_path)
            for i, slide in enumerate(presentation.slides):
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        doc_content += shape.text.strip() + "\n"

        elif file.filename.endswith('.pdf'):
            # Parse the PDF file
            pdf_document = fitz.open(file_path)
            for page_num in range(len(pdf_document)):
                page = pdf_document.load_page(page_num)
                doc_content += page.get_text("text") + "\n"

        else:
            return jsonify({"error": "Invalid file format. Only .pptx and .pdf files are allowed"}), 400

        # Generate summary using OpenAI API
        simplified_content = generate_summary(doc_content)

        # Create DOCX file with simplified content
        doc = Document()
        doc.add_paragraph(simplified_content)
        docx_filename = f"{os.path.splitext(filename)[0]}.docx"
        docx_path = os.path.join(app.config['DOCX_FOLDER'], docx_filename)
        doc.save(docx_path)

        return send_file(docx_path, as_attachment=True)

    except Exception as e:
        return jsonify({"error": f"Error parsing file: {str(e)}"}), 500

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))  # Dynamically set port
    app.run(host='0.0.0.0', port=port)
